"""Google Drive connector for the APFA data pipeline.

Authenticates via a service-account JSON key, recursively lists files in
shared Drive folders, detects changes incrementally, extracts content from
Google Docs (Markdown/TXT), Google Sheets (Sheets API), PDFs (pypdf),
PPTX (python-pptx), and plain text, then emits NormalizedRecord chunks.

Improvements from Perplexity reference:
- Recursive folder listing (sub-folders)
- Chunked downloads via MediaIoBaseDownload (handles large files)
- PPTX / Google Slides support
- Retry decorator on all API calls
- Incremental state tracking (file_id + modifiedTime + content_hash)
- Owner display name extraction from file metadata
- supportsAllDrives for shared drive compatibility

SECURITY:
- Service account auth scoped read-only to shared folders
- Rotate credentials quarterly
- Admin-only — never extend to end users
"""

from __future__ import annotations

import io
import json
import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Generator, Optional

from app.connectors.base import NormalizedRecord, RAGSource
from app.services.chunking import (
    chunk_prose,
    chunk_regulatory,
    section_aware_chunk,
    sentence_aware_chunk,
    spreadsheet_chunk,
)
from app.services.pipeline_utils import (
    load_state,
    now_utc_iso,
    retry,
    save_state,
    sha256_text,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MIME_GOOGLE_DOC = "application/vnd.google-apps.document"
MIME_GOOGLE_SHEET = "application/vnd.google-apps.spreadsheet"
MIME_GOOGLE_SLIDE = "application/vnd.google-apps.presentation"
MIME_PDF = "application/pdf"
MIME_PPTX = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
MIME_DOCX = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
MIME_TXT = "text/plain"
MIME_FOLDER = "application/vnd.google-apps.folder"

SUPPORTED_MIMES = {
    MIME_GOOGLE_DOC,
    MIME_GOOGLE_SHEET,
    MIME_GOOGLE_SLIDE,
    MIME_PDF,
    MIME_PPTX,
    MIME_DOCX,
    MIME_TXT,
}

DRIVE_SCOPES = [
    "https://www.googleapis.com/auth/drive.readonly",
    "https://www.googleapis.com/auth/spreadsheets.readonly",
]

_FILE_FIELDS = (
    "id, name, mimeType, modifiedTime, webViewLink, owners, parents"
)


# ---------------------------------------------------------------------------
# Auth
# ---------------------------------------------------------------------------

def _build_services(
    creds_source: str,
) -> tuple[Any, Any]:
    """Build (drive_service, sheets_service) from service account creds."""
    from google.oauth2 import service_account
    from googleapiclient.discovery import build

    try:
        creds = service_account.Credentials.from_service_account_file(
            creds_source, scopes=DRIVE_SCOPES
        )
    except (FileNotFoundError, ValueError):
        info = json.loads(creds_source)
        creds = service_account.Credentials.from_service_account_info(
            info, scopes=DRIVE_SCOPES
        )

    drive = build("drive", "v3", credentials=creds)
    sheets = build("sheets", "v4", credentials=creds)
    return drive, sheets


# ---------------------------------------------------------------------------
# File listing (recursive)
# ---------------------------------------------------------------------------

def _list_files_recursive(
    drive: Any, folder_id: str
) -> Generator[dict[str, Any], None, None]:
    """Yield all supported files in folder_id and its sub-folders."""
    page_token: Optional[str] = None

    while True:
        query = f"'{folder_id}' in parents and trashed = false"
        resp = (
            drive.files()
            .list(
                q=query,
                fields=f"nextPageToken, files({_FILE_FIELDS})",
                pageSize=200,
                pageToken=page_token,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            )
            .execute()
        )
        for f in resp.get("files", []):
            if f["mimeType"] == MIME_FOLDER:
                yield from _list_files_recursive(drive, f["id"])
            elif f["mimeType"] in SUPPORTED_MIMES:
                yield f

        page_token = resp.get("nextPageToken")
        if not page_token:
            break


# ---------------------------------------------------------------------------
# Change detection
# ---------------------------------------------------------------------------

def _get_start_page_token(drive: Any) -> str:
    """Get current changes page token."""
    resp = drive.changes().getStartPageToken(supportsAllDrives=True).execute()
    return resp["startPageToken"]


def _get_changed_file_ids(
    drive: Any, saved_token: str
) -> tuple[set[str], str]:
    """Return (changed_file_ids, new_page_token) since saved_token."""
    changed: set[str] = set()
    page_token: Optional[str] = saved_token

    while page_token:
        resp = (
            drive.changes()
            .list(
                pageToken=page_token,
                fields="nextPageToken, newStartPageToken, changes(fileId, removed)",
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            )
            .execute()
        )
        for change in resp.get("changes", []):
            changed.add(change["fileId"])

        if "newStartPageToken" in resp:
            return changed, resp["newStartPageToken"]
        page_token = resp.get("nextPageToken")

    return changed, saved_token


# ---------------------------------------------------------------------------
# Content extractors
# ---------------------------------------------------------------------------

@retry(max_attempts=4, base_delay=2.0, exceptions=(Exception,))
def _export_google_doc(drive: Any, file_id: str) -> str:
    """Export a Google Doc as plain text using chunked download."""
    from googleapiclient.http import MediaIoBaseDownload

    request = drive.files().export_media(fileId=file_id, mimeType="text/plain")
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue().decode("utf-8", errors="replace")


@retry(max_attempts=4, base_delay=2.0, exceptions=(Exception,))
def _download_binary(drive: Any, file_id: str) -> bytes:
    """Download a binary file (PDF, PPTX, DOCX) using chunked download."""
    from googleapiclient.http import MediaIoBaseDownload

    request = drive.files().get_media(fileId=file_id, supportsAllDrives=True)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue()


@retry(max_attempts=4, base_delay=2.0, exceptions=(Exception,))
def _get_sheet_data(
    sheets: Any, spreadsheet_id: str
) -> tuple[list[dict[str, Any]], str]:
    """Fetch all sheets from a Google Sheets spreadsheet.

    Returns (list_of_sheet_dicts, spreadsheet_title).
    Each sheet_dict: {name, headers, rows}.
    """
    meta = sheets.spreadsheets().get(spreadsheetId=spreadsheet_id).execute()
    title = meta.get("properties", {}).get("title", "Spreadsheet")
    sheet_metas = meta.get("sheets", [])

    results: list[dict[str, Any]] = []
    for sm in sheet_metas:
        sheet_title = sm["properties"]["title"]
        resp = (
            sheets.spreadsheets()
            .values()
            .get(spreadsheetId=spreadsheet_id, range=sheet_title)
            .execute()
        )
        values: list[list[Any]] = resp.get("values", [])
        if not values:
            continue
        headers = [str(h) for h in values[0]]
        rows = values[1:]
        results.append({"name": sheet_title, "headers": headers, "rows": rows})

    return results, title


def _extract_pdf_text(data: bytes) -> str:
    """Extract text from PDF bytes."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(p for p in pages if p.strip())
    except Exception as exc:
        logger.error("PDF extraction failed: %s", exc)
        return ""


def _extract_docx_text(data: bytes) -> str:
    """Extract text from DOCX bytes."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as exc:
        logger.error("DOCX extraction failed: %s", exc)
        return ""


def _extract_pptx_text(data: bytes) -> str:
    """Extract text from PPTX bytes."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(
                            run.text for run in para.runs if run.text
                        )
                        if line.strip():
                            texts.append(line.strip())
            if texts:
                slides_text.append(
                    f"[Slide {slide_num}]\n" + "\n".join(texts)
                )
        return "\n\n".join(slides_text)
    except ImportError:
        logger.warning("python-pptx not installed — skipping PPTX extraction")
        return ""
    except Exception as exc:
        logger.error("PPTX extraction failed: %s", exc)
        return ""


def _author_from_file(file_meta: dict[str, Any]) -> str:
    """Extract owner display name from file metadata."""
    owners = file_meta.get("owners", [])
    if owners:
        return owners[0].get("displayName", "")
    return ""


# ---------------------------------------------------------------------------
# Incremental state helpers
# ---------------------------------------------------------------------------

def _file_changed(
    file_meta: dict[str, Any], state: dict[str, Any]
) -> bool:
    """Return True if the file's modifiedTime differs from stored state."""
    key = f"file:{file_meta['id']}"
    prev = state.get(key, {})
    return prev.get("modified_time") != file_meta.get("modifiedTime")


def _update_state(
    state: dict[str, Any],
    file_meta: dict[str, Any],
    content_hash: str,
) -> None:
    key = f"file:{file_meta['id']}"
    state[key] = {
        "file_id": file_meta["id"],
        "modified_time": file_meta.get("modifiedTime", ""),
        "content_hash": content_hash,
    }


# ---------------------------------------------------------------------------
# Connector class
# ---------------------------------------------------------------------------

class GoogleDriveConnector(RAGSource):
    """Sync documents from Google Drive into the RAG pipeline.

    Uses a service account with read-only access to shared folders.
    The folder must be explicitly shared with the service account email.
    """

    source_type = "google_drive"

    def __init__(
        self,
        credentials_json: str,
        state_path: str = "/tmp/apfa_drive_state.json",
    ):
        self._creds_source = credentials_json
        self._state_path = Path(state_path)
        self._drive = None
        self._sheets = None

    def _ensure_services(self):
        if self._drive is None:
            self._drive, self._sheets = _build_services(self._creds_source)

    def fetch(
        self,
        folder_ids: list[str] | None = None,
        file_ids: list[str] | None = None,
        force_full: bool = False,
        **kwargs,
    ) -> list[dict]:
        """Fetch documents from Google Drive with incremental change detection.

        Only re-processes files whose modifiedTime changed since last sync.
        Set force_full=True to reprocess everything.
        """
        self._ensure_services()
        state = {} if force_full else load_state(self._state_path)
        raw_docs = []

        if file_ids:
            for fid in file_ids:
                try:
                    meta = (
                        self._drive.files()
                        .get(
                            fileId=fid,
                            fields=_FILE_FIELDS,
                            supportsAllDrives=True,
                        )
                        .execute()
                    )
                    doc = self._extract_file(meta)
                    if doc:
                        raw_docs.append(doc)
                except Exception as e:
                    logger.error(f"Failed to fetch file {fid}: {e}")
            return raw_docs

        if not folder_ids:
            logger.warning("No folder_ids or file_ids provided")
            return []

        for folder_id in folder_ids:
            logger.info(f"Listing files in folder {folder_id}…")
            all_files = list(_list_files_recursive(self._drive, folder_id))
            logger.info(f"Found {len(all_files)} supported files")

            for file_meta in all_files:
                if not force_full and not _file_changed(file_meta, state):
                    logger.debug(
                        "Skipping unchanged: %s", file_meta.get("name", "?")
                    )
                    continue

                doc = self._extract_file(file_meta)
                if doc:
                    raw_docs.append(doc)
                    _update_state(
                        state,
                        file_meta,
                        sha256_text(doc["text"]),
                    )

        save_state(state, self._state_path)
        return raw_docs

    def _extract_file(self, file_meta: dict[str, Any]) -> dict | None:
        """Extract text content from a Drive file."""
        mime = file_meta["mimeType"]
        fid = file_meta["id"]
        name = file_meta.get("name", "Untitled")
        text = ""

        try:
            if mime == MIME_GOOGLE_DOC:
                text = _export_google_doc(self._drive, fid)

            elif mime == MIME_GOOGLE_SHEET:
                sheets_data, ss_title = _get_sheet_data(self._sheets, fid)
                # Store structured data for spreadsheet chunking
                file_meta["_sheets_data"] = sheets_data
                file_meta["_ss_title"] = ss_title
                # Build text representation for transform
                parts = []
                for sheet in sheets_data:
                    header_str = " | ".join(sheet["headers"])
                    parts.append(f"Sheet: {sheet['name']}\nColumns: {header_str}")
                    for row in sheet["rows"][:100]:
                        padded = row + [""] * (len(sheet["headers"]) - len(row))
                        pairs = [
                            f"{h}: {v}"
                            for h, v in zip(sheet["headers"], padded)
                            if str(v).strip()
                        ]
                        parts.append("; ".join(pairs))
                text = "\n".join(parts)

            elif mime == MIME_PDF:
                data = _download_binary(self._drive, fid)
                text = _extract_pdf_text(data)

            elif mime == MIME_DOCX:
                data = _download_binary(self._drive, fid)
                text = _extract_docx_text(data)

            elif mime in (MIME_PPTX, MIME_GOOGLE_SLIDE):
                if mime == MIME_GOOGLE_SLIDE:
                    from googleapiclient.http import MediaIoBaseDownload

                    request = self._drive.files().export_media(
                        fileId=fid, mimeType=MIME_PPTX
                    )
                    buf = io.BytesIO()
                    downloader = MediaIoBaseDownload(buf, request)
                    done = False
                    while not done:
                        _, done = downloader.next_chunk()
                    data = buf.getvalue()
                else:
                    data = _download_binary(self._drive, fid)
                text = _extract_pptx_text(data)

            elif mime == MIME_TXT:
                data = _download_binary(self._drive, fid)
                text = data.decode("utf-8", errors="replace")

        except Exception as e:
            logger.error(f"Text extraction failed for {name} ({fid}): {e}")
            return None

        if not text or not text.strip():
            logger.warning(f"Empty content extracted from {name}")
            return None

        file_meta["text"] = text.strip()
        file_meta["_author"] = _author_from_file(file_meta)
        return file_meta

    def transform(self, raw_docs: list[dict]) -> list[NormalizedRecord]:
        """Transform raw Drive documents into chunked NormalizedRecords."""
        now_iso = now_utc_iso()
        all_records: list[NormalizedRecord] = []

        for doc in raw_docs:
            doc_id = f"gdrive:{doc['id']}"
            mime = doc.get("mimeType", "")
            text = doc["text"]
            author = doc.get("_author", "")

            # Choose chunking strategy
            if mime == MIME_GOOGLE_SHEET and "_sheets_data" in doc:
                # Use structured spreadsheet chunking
                chunk_dicts = []
                for sheet in doc["_sheets_data"]:
                    chunk_dicts.extend(
                        spreadsheet_chunk(
                            sheet_name=sheet["name"],
                            headers=sheet["headers"],
                            rows=sheet["rows"],
                        )
                    )
            elif self._looks_regulatory(text):
                chunk_dicts = section_aware_chunk(text)
            else:
                chunk_dicts = sentence_aware_chunk(text)

            for chunk in chunk_dicts:
                idx = chunk["chunk_index"]
                ext_id = f"google_drive:{doc['id']}:{idx:04d}"
                content_kind = "doc_section"
                if mime == MIME_GOOGLE_SHEET:
                    content_kind = "spreadsheet"
                elif mime in (MIME_PPTX, MIME_GOOGLE_SLIDE):
                    content_kind = "presentation"

                record = NormalizedRecord(
                    external_id=ext_id,
                    source_type="google_drive",
                    source_url=doc.get(
                        "webViewLink",
                        f"https://drive.google.com/file/d/{doc['id']}",
                    ),
                    title=doc.get("_ss_title", doc.get("name", "")),
                    author=author,
                    published_at=doc.get("modifiedTime", now_iso),
                    fetched_at=now_iso,
                    freshness_class="static",
                    content_kind=content_kind,
                    text=chunk["text"],
                    chunk_index=idx,
                    total_chunks=len(chunk_dicts),
                    parent_document_id=doc_id,
                    metadata_json=json.dumps(
                        {
                            "drive_file_id": doc["id"],
                            "mime_type": mime,
                            "section_title": chunk.get("section_title", ""),
                            "chunk_kind": chunk.get("chunk_kind", ""),
                        }
                    ),
                )
                all_records.append(record)

        return all_records

    @staticmethod
    def _looks_regulatory(text: str) -> bool:
        """Heuristic: does this text look like a regulatory document?"""
        indicators = [
            "Section ",
            "SECTION ",
            "§",
            "pursuant to",
            "regulation",
            "compliance",
            "disclosure",
            "Act of ",
            "Code of Federal Regulations",
        ]
        score = sum(1 for ind in indicators if ind in text[:2000])
        return score >= 3

    def get_changes_token(self) -> str:
        """Get current changes page token for incremental sync."""
        self._ensure_services()
        return _get_start_page_token(self._drive)
